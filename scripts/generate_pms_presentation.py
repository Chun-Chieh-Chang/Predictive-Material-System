#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PMS 跨部門系統操作與決策協同培訓簡報生成器 (18 頁超詳細 + 4 大情境實戰模擬版)
利用 python-pptx 生成 16:9 寬螢幕、Bento 卡片式佈局、嚴格色彩規範與完整逐字級演講備忘錄 (Speaker Notes)。
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 設計色彩規範 (Design Tokens) ──────────────────────────────────────
C_NAVY_DARK    = RGBColor(15, 23, 42)     # #0F172A 深 Slate 底色/主標題
C_BLUE_BRAND   = RGBColor(14, 116, 144)   # #0E7490 品牌深青/鈷藍
C_BLUE_ACCENT  = RGBColor(2, 132, 199)    # #0284C7 亮藍 Accent
C_EMERALD      = RGBColor(16, 185, 129)   # #10B981 成功/通過
C_AMBER        = RGBColor(245, 158, 11)   # #F59E0B 警示/注意
C_RED          = RGBColor(239, 68, 68)    # #EF4444 錯誤/缺料
C_PURPLE       = RGBColor(147, 51, 234)   # #9333EA 特殊指標/模擬
C_BG_LIGHT     = RGBColor(248, 250, 252)  # #F8FAFC 淺灰白卡片底色
C_CARD_BORDER  = RGBColor(226, 232, 240)  # #E2E8F0 卡片邊框
C_TEXT_MAIN    = RGBColor(30, 41, 59)     # #1E293B 主文字
C_TEXT_MUTED   = RGBColor(100, 116, 139)  # #64748B 次要說明文字
C_WHITE        = RGBColor(255, 255, 255)  # 純白
C_CARD_HEADER  = RGBColor(241, 245, 249)  # #F1F5F9 淺標題底

FONT_TITLE = "微軟正黑體"
FONT_BODY  = "微軟正黑體"

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_header(slide, tag, title, subtitle):
        """頂部標題列"""
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.45), Inches(0.18), Inches(0.85))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C_BLUE_ACCENT
        bar.line.fill.background()

        tx_box = slide.shapes.add_textbox(Inches(1.08), Inches(0.40), Inches(11.5), Inches(0.95))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p1 = tf.paragraphs[0]
        p1.text = f"【{tag}】 {title}"
        p1.font.name = FONT_TITLE
        p1.font.size = Pt(21)
        p1.font.bold = True
        p1.font.color.rgb = C_NAVY_DARK

        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = FONT_BODY
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.space_before = Pt(3)

    def add_card(slide, left, top, width, height, title, items, badge_text="", badge_color=C_BLUE_ACCENT, bg_color=C_BG_LIGHT):
        """標準卡片組件"""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = C_CARD_BORDER
        shape.line.width = Pt(1.5)

        tx = slide.shapes.add_textbox(Inches(left + 0.22), Inches(top + 0.18), Inches(width - 0.44), Inches(height - 0.36))
        tf = tx.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(13.5)
        p_title.font.bold = True
        p_title.font.color.rgb = C_NAVY_DARK

        if badge_text:
            p_badge = tf.add_paragraph()
            p_badge.text = f"▶ {badge_text}"
            p_badge.font.name = FONT_BODY
            p_badge.font.size = Pt(9.5)
            p_badge.font.bold = True
            p_badge.font.color.rgb = badge_color
            p_badge.space_after = Pt(4)

        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.name = FONT_BODY
            p.font.size = Pt(10.0)
            p.font.color.rgb = C_TEXT_MAIN
            p.space_after = Pt(3)

    # ═════════════════════════════════════════════════════════════════════
    # Slide 1: 封面 (Title Slide)
    # ═════════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_NAVY_DARK
    bg1.line.fill.background()

    halo = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    halo.fill.solid()
    halo.fill.fore_color.rgb = RGBColor(30, 41, 59)
    halo.line.color.rgb = C_BLUE_ACCENT
    halo.line.width = Pt(2)

    tx1 = s1.shapes.add_textbox(Inches(1.4), Inches(1.3), Inches(10.5), Inches(5.0))
    tf1 = tx1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "QCC 料事如神圈 · 射出成型智能備料與排程推估平台"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_BLUE_ACCENT
    p.space_after = Pt(10)

    p = tf1.add_paragraph()
    p.text = "料事如神系統 (PMS)\n跨部門操作實務、決策協同與情境演練手冊"
    p.font.name = FONT_TITLE
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.space_after = Pt(14)

    p = tf1.add_paragraph()
    p.text = "消除供需時序落差 · 三階 MRP 自動展開 · 4 大實戰情境動態推演 · 支援週二出貨排程精準放行"
    p.font.name = FONT_BODY
    p.font.size = Pt(12.5)
    p.font.color.rgb = RGBColor(148, 163, 184)
    p.space_after = Pt(24)

    p = tf1.add_paragraph()
    p.text = "報告對象：業務部 · 製造部 · 工程部 · 品保部 · 資材部 | 主講人：Wesley Chang | 2026 年 8 月"
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = C_EMERALD

    s1.notes_slide.notes_text_frame.text = (
        "【演講備忘錄 - Slide 1】：各位長官、各部門主管以及參與 QCC 的專案同仁大家好。\n"
        "今天由我代表『料事如神圈』向大家正式匯報『PMS 智能備料與排程推估系統』的跨部門操作實務與實戰情境演練。\n"
        "今天簡報除了說明各部門職責與介面外，我們特別規劃了『4 大工廠最常發生的突發情境模擬』，"
        "讓大家親眼見證系統如何幫助工廠在 5 分鐘內做出最佳排程與備料決策。"
    )

    # ═════════════════════════════════════════════════════════════════════
    # Slide 2 ~ Slide 10: 基礎理論與部門操作手冊
    # ═════════════════════════════════════════════════════════════════════
    # Slide 2: 現狀剖析
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "現狀剖析", "製造與供應鏈三大核心痛點深度剖析", "資訊不對稱與長交期原料引發的經營與排程風險")
    add_card(s2, 0.8, 1.55, 3.6, 5.35, "痛點一：預測波動劇烈", [
        "ICU 客戶行銷端提供未來半年/一年 Forecast，但歷史準確率僅 35%~41%，波動極大。",
        "過去因缺乏客觀比對工具，業務端在每週視訊會議中無法有效反駁客戶的突然砍單或追加。",
        "缺乏我方為其備料的客觀數據背書，商務談判處於被動責難局面。"
    ], "商務談判痛點", C_RED)
    add_card(s2, 4.8, 1.55, 3.6, 5.35, "痛點二：週二協調耗時", [
        "每兩週接收 ICU Ship Schedule，業務需在每週二開會決議實際承接量。",
        "業務與生管手動拉 Excel 表比對耗時 2 小時，容易因公式人為疏漏造成算錯。",
        "夜班 12h 無人挑選之在製品 (WIP) 存在時序差，白天開會常誤判缺料不敢放行。"
    ], "現場排程痛點", C_AMBER)
    add_card(s2, 8.8, 1.55, 3.6, 5.35, "痛點三：爆倉與缺料兩難", [
        "德國 INEOS (MABS) 與美國 Avient (PVC) 原料採購交期長達 3~4 個月，MOQ 高達 5 噸。",
        "客戶需求驟降時採購端無法即時煞車，引發高達 8,000 萬元呆滯庫存風險。",
        "原料倉實體容積上限僅 10 噸，一次進櫃極易塞滿走道引發安全疑慮。"
    ], "庫存資金痛點", C_RED)
    s2.notes_slide.notes_text_frame.text = "【演講備忘錄 - Slide 2】：這三大痛點是工廠日常最痛的死結，PMS 系統正是為了解決這三大痛點而設計。"

    # Slide 3: 系統總體架構
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "架構總覽", "PMS 系統全域架構與單一事實來源 (SSOT)", "4 大運算引擎 · 10 大標準主檔 · 5 層物料分類體系")
    add_card(s3, 0.8, 1.55, 5.6, 2.5, "全廠單一事實來源 (SSOT)", [
        "徹底打破部門資訊孤島，消除各自維護不同版本 Excel 的混亂。",
        "將業務預測、實際訂單、機台模具、在庫與 WIP 整合至同一資料庫。",
        "提供 10 大標準主檔維護與 3 級變更管制，確保數據高一致性。"
    ], "資料層 (Data Layer)", C_BLUE_BRAND)
    add_card(s3, 6.8, 1.55, 5.6, 2.5, "五層階層化物料分類 (MECE)", [
        "RAW (原料類，KG)：塑膠原粒、色母、色粉，受供應商規則管制。",
        "MAT (物料類，PCS)：滅菌袋、紙箱、標籤、乾燥劑等包裝輔料。",
        "PART (零件類，PCS)：單品射出成型件 (多數為直接出貨品)。",
        "COMP (組件類，PCS)：次總成裝配件；SET (套件類)：最終組合成品。"
    ], "領域模型 (Domain Model)", C_BLUE_ACCENT)
    add_card(s3, 0.8, 4.35, 5.6, 2.55, "4 大高精度運算與推導引擎", [
        "MRP 三階展開計算引擎：需求淨額 → BOM爆炸 → 採購建議與分批到貨。",
        "WIP 日動態推估公式計算器：日累積模型 + 夜班 12h 時序差補償。",
        "訂單全鏈路物料緊張度診斷引擎：6 大環節掃描 + 4 級色標告警。",
        "全數據鏈路深度模擬器：4 大業務場景穿透與孤兒數據排查。"
    ], "運算層 (Computing Engines)", C_EMERALD)
    add_card(s3, 6.8, 4.35, 5.6, 2.55, "智慧雙模換檔與資料中心", [
        "開箱預載 52 筆代表性全階層示範演練庫 (DEMO 模式)，零冷啟動。",
        "匯入真實 Excel / JSON 時自動無縫換檔為【🟢 正式生產模式】。",
        "支援隨時一鍵重載示範包或一鍵清空，供日常培訓與正式維運切換。"
    ], "應用層 (Application Gateway)", C_AMBER)
    s3.notes_slide.notes_text_frame.text = "【演講備忘錄 - Slide 3】：系統架構統一了資料庫與演算法，全廠擁有唯一的運算基準。"

    # Slide 4: 3階 MRP 演算法
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "核心演算法", "3 階 MRP 推導與倒推發單排程演算法", "精確數學模型 · 消除 8,000 萬呆料與防爆倉分批進貨")
    add_card(s4, 0.8, 1.55, 3.6, 5.35, "步驟 1：成品淨需求計算", [
        "需求量來源：依需求沖銷模式 (疊加/沖銷/僅預估/僅實單) 計算總需求。",
        "成品淨需求 = 總需求 - 成品良品現貨 - 待驗品(WIP×全檢良率)。",
        "若成品現有量足夠，則淨需求為 0，無需啟動射出排程。"
    ], "Step 1: Net Requirement", C_BLUE_ACCENT)
    add_card(s4, 4.8, 1.55, 3.6, 5.35, "步驟 2：BOM 展開與損耗折算", [
        "單件製品原料毛需求 = (單模淨重 + 料頭流道重/穴數) / 1000 (KG)。",
        "考慮成型標準損耗率：實際原料需用量 = 毛需求 / (1 - 損耗率)。",
        "考慮全檢良率：總原料毛需求 = 實際原料需用量 / 全檢良率標準。",
        "若有色母配色需求，依添加比例 (如 2%) 同步展開色母原料毛需求。"
    ], "Step 2: BOM Explosion", C_PURPLE)
    add_card(s4, 8.8, 1.55, 3.6, 5.35, "步驟 3：倒推採購日與分批進貨", [
        "原料淨需求 = 總原料毛需求 - 原料在庫 - 在途採購PO + 虛擬預扣量。",
        "MOQ 取整：採購量 = 向上取整至供應商 MOQ (如 5 噸) 的整倍數。",
        "倒推下單期限：Order Deadline = 需求交期 - 原料採購交期 - 製造工時天數。",
        "分批進貨防爆倉：若採購量 > 倉容上限 (10噸)，自動拆分為多批進櫃。"
    ], "Step 3: Phased Procurement", C_EMERALD)
    s4.notes_slide.notes_text_frame.text = "【演講備忘錄 - Slide 4】：MRP 演算法把成品缺口一層層拆解到塑膠粒與色母，並精準計算下單日期與進貨批次。"

    # Slide 5: WIP 推估演算法
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "核心演算法", "在製品 (WIP) 日動態推估與夜班時序差補償", "日累積連續推演模型 · 消除 12 小時現場無人挑選時序差")
    add_card(s5, 0.8, 1.55, 3.6, 5.35, "WIP 日累積動態模型", [
        "數學公式：WIP(t) = WIP(t-1) + P(t) - S(t)",
        "WIP(t-1)：前一日結存之車間待驗在製品數量。",
        "P(t)：當日射出機台理論產出總量。",
        "S(t)：當日 3F 全檢作業實際挑選完成並入庫之良品數量。"
    ], "動態公式推導", C_BLUE_ACCENT)
    add_card(s5, 4.8, 1.55, 3.6, 5.35, "機台理論產出 P(t) 計算", [
        "P(t) = 運轉工時(h) × (3600 / 週期) × 妥善穴數 × (1 - 損耗率)",
        "例如：運轉 24 小時、週期 25 秒、16 穴模、損耗 3%：",
        "P(t) = 24 × (3600/25) × 16 × 0.97 = 53,683 PCS/日。",
        "若模具有塞穴，自動以妥善穴數 (如 14 穴) 重新計算。"
    ], "機台產能公式", C_EMERALD)
    add_card(s5, 8.8, 1.55, 3.6, 5.35, "夜班 12h 時序差補償機制", [
        "現場痛點：夜班 12 小時射出機持續運轉，但品保全檢人員未上班，產出堆積在 3F 待驗區。",
        "隔天週二上午開會時，舊系統帳面良品庫存顯示為 0，容易誤判實質缺貨。",
        "PMS 自動補償：有效可用量 = 成品現貨 + (WIP + 夜班推估產出) × 全檢良率標準。"
    ], "時序差消除", C_AMBER)
    s5.notes_slide.notes_text_frame.text = "【演講備忘錄 - Slide 5】：WIP 引擎解決了夜班生產但沒全檢的帳面盲區，確保週二決策依據是最真實的庫存。"

    # Slide 6: 訂單緊張度診斷
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "核心演算法", "訂單全鏈路物料緊張度雷達診斷引擎", "6 大供應鏈環節瓶頸掃描 · 4 級緊張度色標與應變處置 SOP")
    add_card(s6, 0.8, 1.55, 3.6, 5.35, "6 大供應鏈瓶頸掃描環節", [
        "1. 成品在庫水位 (FG Stock)：現貨是否可 100% 覆蓋訂單？",
        "2. 待驗 WIP 折算 (WIP Capacity)：3F 待驗品挑選後是否足夠？",
        "3. 機台日產能極限 (Daily Capacity)：製造交期是否趕得上？",
        "4. 原料在庫充足度 (RM On-hand)：廠內原料是否足夠射出？",
        "5. 在途採購進度 (PO In-transit)：海運/報關/到廠 ETA 是否逾期？",
        "6. 供應商交期緩衝 (Lead Time Buffer)：是否已進入安全交期紅線？"
    ], "全鏈路 6 大節點", C_BLUE_ACCENT)
    add_card(s6, 4.8, 1.55, 3.6, 5.35, "4 級緊張度色標評級標準", [
        "🔴 致命缺料 (Critical Shortage)：原料或產能實質赤字，勢必延誤出貨。",
        "🟡 產能緊張 (Capacity Tension)：現有良品不足，需 3F WIP 優先挑選支援。",
        "🔵 即刻發單 (PO Urgent)：已進入供應商最晚發單期限 (Order Deadline)。",
        "🟢 安全充裕 (Safe Buffer)：庫存、WIP 與在途原料全量覆蓋需求。"
    ], "4 級告警標準", C_RED)
    add_card(s6, 8.8, 1.55, 3.6, 5.35, "逐筆訂單根因分析與應變 SOP", [
        "系統提供逐筆訂單點擊展開功能，自動輸出根本原因分析 (RCA)。",
        "配套標準處置 SOP：如提示生管『啟動 3F 夜班挑選加班』或提示採購『催促海運報關加速』。",
        "支援全文檢索、客戶過濾與緊急程度排序，一目了然。"
    ], "處置 SOP 導引", C_EMERALD)
    s6.notes_slide.notes_text_frame.text = "【演講備忘錄 - Slide 6】：緊張度看板像雷達一樣隨時掃描每張訂單會卡在哪個節點，並給出標準 SOP。"

    # Slide 7: 業務實戰
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "部門實戰", "業務部門：週二出貨排程審查看板操作實務", "每週二出貨協調會議專用 · 5 分鐘完成雙週排程放行決策")
    add_card(s7, 0.8, 1.55, 3.6, 5.35, "會議前：資料準備與載入", [
        "每兩週週一接收 ICU 或 MDX 最新出貨排程 (Ship Schedule)。",
        "至『數據中心 ➔ 10大主檔維護』更新『業務預估需求檔』或至『無損資料中心』匯入 Excel。",
        "確認品號與客戶料號對應正確 (支援 A01-200-131 與 R1-2355 雙品號)。"
    ], "Step 1: 資料就緒", C_BLUE_ACCENT)
    add_card(s7, 4.8, 1.55, 3.6, 5.35, "會議中：三色燈號審查放行", [
        "打開『週二出貨審查看板 (ShipScheduleClearanceView)』。",
        "🟢 綠燈 (100% 可放行)：成品良品現貨充足，直接回覆客戶可準時出貨。",
        "🟡 黃燈 (需 WIP 支援)：現貨不足但 3F 待驗品充足，現場指示生管優先挑選排程。",
        "🔴 紅燈 (實質缺貨)：現貨與待驗品皆不足，現場協調插單生產或與客戶協商交期。"
    ], "Step 2: 燈號審查", C_EMERALD)
    add_card(s7, 8.8, 1.55, 3.6, 5.35, "What-If 需求動態模擬滑桿", [
        "拖曳看板頂部的『What-If 需求模擬滑桿 (0.5x ~ 2.0x)』。",
        "即時模擬：如果 ICU 突然追加 50% 訂單，廠內產能與 WIP 是否接得住？",
        "模擬結果即時動態連動三色燈號與缺口數量，輔助高階商務決策。"
    ], "Step 3: What-If 模擬", C_PURPLE)
    s7.notes_slide.notes_text_frame.text = "【演講備忘錄 - Slide 7】：業務同仁每週二開會直接看燈號做決定，拉動 What-If 滑桿就能預測客戶追加訂單的承接能力。"

    # Slide 8: 業務商務談判
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "部門實戰", "業務部門：預測偏差分析與客戶供需透明化報告", "決策戰情室商務賦能 · 消除跨國客戶客訴與商務談判被動")
    add_card(s8, 0.8, 1.55, 3.6, 5.35, "歷史預測偏差率計算", [
        "數學公式：Deviation % = (Actual PO - Forecast) / Forecast × 100%",
        "系統自動比對歷史各週次之預測需求與正式下單實績。",
        "預測準確度評級條：紅色 (<50% 極度不穩定)、黃色 (50~80% 尚可)、綠色 (>80% 高度準確)。"
    ], "預測偏差客觀量化", C_RED)
    add_card(s8, 4.8, 1.55, 3.6, 5.35, "供需透明化備料背書", [
        "即時展示我方為該客戶因應 Forecast 所準備的特定原料總量。",
        "客觀佐證結構：原料需求量 vs (廠內在庫良品 + 在途海運進櫃量)。",
        "視覺化展示備料覆蓋率，消除客戶對我方『準備不足、資訊不透明』之質疑。"
    ], "備料客觀透明度", C_EMERALD)
    add_card(s8, 8.8, 1.55, 3.6, 5.35, "跨國視訊商務談判話術應用", [
        "當 ICU 採購主管質疑交期時，直接出示戰情室報表。",
        "以客觀數據說明：『貴司今年某產品預估 100 萬，實下 50 萬，但我方仍依約備妥 8 噸原料在途』。",
        "爭取更有利的 MOQ 承諾、長交期原料共擔責任或更長的出貨寬限期。"
    ], "商務談判背書", C_BLUE_BRAND)
    s8.notes_slide.notes_text_frame.text = "【演講備忘錄 - Slide 8】：供需透明化報告是業務跟國外客戶開會最強大的談判武器。"

    # Slide 9: 製造生管實戰
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "部門實戰", "製造與生管：模具狀態維護與車間產能排程", "責任角色：製造主管 / 車間組長 / 物料生管人員")
    add_card(s9, 0.8, 1.55, 3.6, 5.35, "模具主檔動態維護 (mold_master)", [
        "15 組主力射出模具完整建檔 (涵蓋 2 穴至 32 穴)。",
        "狀態管理：active(常態量產)、trial(試模)、maintenance(維修保養/修穴)、retired(封存報廢)。",
        "動態塞穴調整：若 24 穴模具塞 2 穴，將妥善穴數改為 22，系統產能與原料耗用立刻自動聯動重新計算！"
    ], "模具與穴數管制", C_BLUE_ACCENT)
    add_card(s9, 4.8, 1.55, 3.6, 5.35, "車間產能與成型週期掌控", [
        "精準維護成型週期 (Cycle Time，單位：秒)，如 T接頭 25.0s、直通接頭 18.0s、針筒外筒 32.0s。",
        "系統自動推算單日極限產能與換模工時，防止生管超排機台負荷。",
        "支援主模與備用模切換 (如 16 穴主模與 8 穴備用模自由選擇)。"
    ], "機台產能規劃", C_EMERALD)
    add_card(s9, 8.8, 1.55, 3.6, 5.35, "月內自用料虛擬預扣 (Virtual Backflush)", [
        "解決傳統頂新 ERP 月底才扣料導致月中可用庫存虛增之盲區。",
        "在『系統參數配置』中可啟用『虛擬預扣開關』。",
        "排程展開時自動預先扣除月中已排產的自用料，防止因帳面虛假庫存而少訂料停線。"
    ], "防呆扣料機制", C_AMBER)
    s9.notes_slide.notes_text_frame.text = "【演講備忘錄 - Slide 9】：製造生管同仁只要維護好模具穴數與週期，全廠排程自動連動計算。"

    # Slide 10: 工程品保實戰
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "部門實戰", "工程與品保：產品模具 BOM 與全檢良率標準", "責任角色：工程主管 / 機構工程師 / 品保主管 / 檢驗員")
    add_card(s10, 0.8, 1.55, 3.6, 5.35, "產品模具 BOM (product_mold_bom)", [
        "品號關聯模具代碼、主力原料品號 (如 RAW-ABS-2802 / RAW-PC-110)。",
        "精準填報單模淨重 (net_mold_weight_g) 與料頭流道重 (runner_weight_g)。",
        "設定標準成型損耗率 (std_mfg_scrap_rate，如 2%~5%) 與色母添加比 (如 2.0%)。"
    ], "BOM 核心參數", C_BLUE_BRAND)
    add_card(s10, 4.8, 1.55, 3.6, 5.35, "全檢良率標準 (yield_master)", [
        "維護 3F 車間全檢合格率基準 (如 98.5% / 96.0%)。",
        "作為 3 階 MRP 推導時折算原料毛需求之底層安全係數。",
        "作為週二出貨審查時折算 WIP 有效可用量之乘數，防止不良品流入出貨計算。"
    ], "品質標準基準", C_EMERALD)
    add_card(s10, 8.8, 1.55, 3.6, 5.35, "品質紀錄與批次追溯閉環", [
        "全檢實際良率日報 (sorting_actual_yield_log)：記錄每日批號、全檢數、合格數與實際良率。",
        "色母/色粉混合日報 (color_mixing_log)：記錄基礎樹脂與色母配比批次，實現醫療級 UDI 完整追溯。",
        "良率異常主動回饋機制。"
    ], "閉環與追溯", C_BLUE_ACCENT)
    s10.notes_slide.notes_text_frame.text = "【演講備忘錄 - Slide 10】：工程與品保維護的克重、損耗與良率是系統最核心的運算基石。"

    # ═════════════════════════════════════════════════════════════════════
    # Slide 11 ~ Slide 14: 4 大實戰情境動態模擬演練 (Scenario Simulations)
    # ═════════════════════════════════════════════════════════════════════

    # Slide 11: 情境模擬一 — 黑天鵝急單插單 / 爆量追加 (+50%)
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "情境演練 1", "情境模擬一：黑天鵝急單插單 / 客戶爆量追加 (+50%)", "實戰模擬：ICU 週二前夕突然追加 50,000 PCS T接頭出貨需求")

    add_card(s11, 0.8, 1.55, 3.6, 5.35, "1. 突發業務情境 (Scenario)", [
        "週二上午 09:00 出貨協調會前，ICU 採購突發特急信件：",
        "T接頭 (A01-200-131) 本週出貨需求由原訂 100,000 PCS 緊急拉高至 150,000 PCS (+50%)！",
        "傳統做法：業務與生管手動翻 Excel 表計算 2 小時，互相爭執到底交不交得出來。"
    ], "情境觸發", C_RED)

    add_card(s11, 4.8, 1.55, 3.6, 5.35, "2. PMS 系統操作演示 (Action)", [
        "1. 打開『週二出貨審查看板』，拖曳 What-If 滑桿至 1.5x (150%)。",
        "2. 系統瞬間重新運算：放行燈號由 🟢 綠燈轉為 🟡 黃燈 (需 WIP 支援)。",
        "3. 數據透視：現有良品庫存 15,000 PCS，3F 待驗區有 WIP 25,000 PCS，夜班已產出 20,000 PCS。",
        "4. 切換至『3階 MRP 推導器』：原料 MABS 在途已有 5,000 KG 預計下週到港。"
    ], "系統即時響應", C_PURPLE)

    add_card(s11, 8.8, 1.55, 3.6, 5.35, "3. 跨部門協同處置 (Decision)", [
        "生管決策：當天立即安排 3F 全檢人員優先挑選 A01-200-131 批次，預計今日可入庫 44,000 PCS。",
        "採購決策：確認在途原料足夠支撐下週持續射出，無需緊急空運原料。",
        "業務決策：09:05 會議中直接回覆客戶：『15 萬全數接單，分兩批於本週五與下週二完整出貨！』"
    ], "5分鐘決策放行", C_EMERALD)

    s11.notes_slide.notes_text_frame.text = (
        "【演講備忘錄 - Slide 11】：這是我們工廠最常碰到的『黑天鵝急單追加』情境！\n"
        "客戶突然多要 5 萬顆，過去大家在會議室吵半天。現在業務只要把滑桿拉到 1.5x，\n"
        "系統立刻顯示黃燈，告訴生管『待驗區有貨，只要今天挑出來就能交』！\n"
        "採購同步確認原料在途充足，前後不到 5 分鐘，業務就能自信滿滿地答應客戶接單！"
    )

    # Slide 12: 情境模擬二 — 客戶突發砍單下修 / 採購踩煞車防爆倉
    s12 = prs.slides.add_slide(blank_layout)
    add_header(s12, "情境演練 2", "情境模擬二：客戶突發預測下修 / 採購踩煞車防 8,000 萬爆倉", "實戰模擬：MDX 突發砍單 30,000 PCS，系統精準攔截過量採購")

    add_card(s12, 0.8, 1.55, 3.6, 5.35, "1. 突發業務情境 (Scenario)", [
        "MDX 因歐美終端通路銷售疲軟，週一更新 Forecast：",
        "將 11 月 Y管 (C09-200-251) 預估需求由 120,000 PCS 驟降至 90,000 PCS (-25%)！",
        "傳統做法：採購端往往不知情，仍依舊版預測向德國下單 5 噸 MABS 原料，導致年底爆倉積壓 8,000 萬元呆滯料！"
    ], "情境觸發", C_RED)

    add_card(s12, 4.8, 1.55, 3.6, 5.35, "2. PMS 系統操作演示 (Action)", [
        "1. 業務在主檔匯入最新預測，系統在『決策戰情室』即時觸發 🟡 預測下修防爆倉警報。",
        "2. 採購打開『3階 MRP 推導器』執行展開：",
        "3. 系統比對：成品淨缺口縮小，計算出原料毛需求僅剩 2,250 KG。",
        "4. 扣除廠內現有庫存 2,450 KG 後，原料淨需求轉為 0 KG！建議採購量自動由 5,000 KG 歸零！"
    ], "系統即時響應", C_AMBER)

    add_card(s12, 8.8, 1.55, 3.6, 5.35, "3. 跨部門協同處置 (Decision)", [
        "採購處置：系統 1 秒內踩住煞車，暫停向德國原廠發出 5 噸 (約 2 萬美元) 的 PO 發單作業。",
        "倉管處置：避免 5 噸大宗原料進廠塞爆原料倉 (上限 10 噸)。",
        "管理成效：成功消除數百萬元營運資金積壓，徹底解決年終呆滯存貨問題！"
    ], "資金防禦成效", C_EMERALD)

    s12.notes_slide.notes_text_frame.text = (
        "【演講備忘錄 - Slide 12】：情境二剛好相反，是『客戶砍單踩煞車』。\n"
        "過去業務收到客戶下修預測，採購不知道，照樣買了 5 噸德國原料，等貨漂洋過海到港才發現沒人要用，"
        "年底倉庫被塞爆，幾百萬資金就壓在那裡。\n"
        "現在 PMS 系統只要業務一把預測改小，採購打開 MRP 看到建議採購量直接變成 0，"
        "一秒鐘踩住煞車，替公司守住珍貴的現金流！"
    )

    # Slide 13: 情境模擬三 — 車間模具突發咬模塞穴 (16穴降14穴)
    s13 = prs.slides.add_slide(blank_layout)
    add_header(s13, "情境演練 3", "情境模擬三：車間模具突發咬模塞穴 / 備用模雙線動態調配", "實戰模擬：主力 16 穴模咬模塞 2 穴，系統動態重算產能與原料單耗")

    add_card(s13, 0.8, 1.55, 3.6, 5.35, "1. 突發車間情境 (Scenario)", [
        "週三上午 10:30，1號廠射出機 A-03 上的主力 16 穴 T接頭模具 (MI-T-16C) 發生第 7、19 穴咬模！",
        "現場組長緊急塞穴處理，實際稼動穴數降為 14 穴。",
        "傳統做法：生管不知道塞穴，仍按 16 穴排產，導致月底交貨短缺 15%，且原料單耗計算失真。"
    ], "情境觸發", C_AMBER)

    add_card(s13, 4.8, 1.55, 3.6, 5.35, "2. PMS 系統操作演示 (Action)", [
        "1. 製造主管打開『10大主檔維護 ➔ 模具主檔』，將妥善穴數由 16 改為 14。",
        "2. 系統後台自動聯動：日產能瞬間由 53,683 PCS 重算為 46,972 PCS (-12.5%)。",
        "3. 單模流道料頭重量分攤自動重算：單件製品毛耗由 1.12g 微調至 1.21g。",
        "4. 『訂單緊張度追蹤看板』立即發出 🟡 產能緊張預警，提示出貨寬限期壓縮 2.5 天！"
    ], "系統即時響應", C_PURPLE)

    add_card(s13, 8.8, 1.55, 3.6, 5.35, "3. 跨部門協同處置 (Decision)", [
        "生管處置：生管立即啟用備用 8 穴模 (MI-T-08C) 於射出機 A-07 進行雙線併行排產，補齊 12.5% 產能缺口。",
        "工程處置：工程部開立修模單，排定週五夜間進行模具超音波清洗與修穴。",
        "結果：交期 100% 履約，原料耗用精準扣帳，無任何帳料偏差。"
    ], "雙線動態平衡", C_EMERALD)

    s13.notes_slide.notes_text_frame.text = (
        "【演講備忘錄 - Slide 13】：情境三是車間最常發生的『模具塞穴』。\n"
        "16 穴模具塞了 2 穴，現場組長只要在系統把穴數改成 14，\n"
        "系統立刻告訴生管『產能掉了 12.5%，交期會被壓縮』！生管馬上把備用的 8 穴模開起來雙線生產，"
        "工程部排修模，客戶完全感受不到任何延遲，這就是真正的敏捷製造！"
    )

    # Slide 14: 情境模擬四 — 國外在途海運延誤 25 天
    s14 = prs.slides.add_slide(blank_layout)
    add_header(s14, "情境演練 4", "情境模擬四：國外在途海運延誤 25 天 / 提前借料與調產預警", "實戰模擬：德國原料海運 ETA 延期，系統提早 45 天發出斷料警報")

    add_card(s14, 0.8, 1.55, 3.6, 5.35, "1. 突發物流情境 (Scenario)", [
        "德國 INEOS 航運公司發出通知：5 噸 MABS 原料 (PO-RM-2026-0501) 因紅海繞道與海關查驗，",
        "預計到港日 (ETA) 由原訂 9/15 嚴重延誤至 10/10 (+25 天)！",
        "傳統做法：等到 9/15 倉庫沒收到貨才發現斷料，射出機台被迫停線，引發客戶嚴重客訴與罰款！"
    ], "情境觸發", C_RED)

    add_card(s14, 4.8, 1.55, 3.6, 5.35, "2. PMS 系統操作演示 (Action)", [
        "1. 資材人員在『在途訂單檔』將狀態改為 delayed 並將 ETA 更新為 2026-10-10。",
        "2. 『訂單緊張度追蹤看板』即時雷達掃描：針對 10 月初交期之訂單亮起 🔴 致命斷料紅燈！",
        "3. 系統自動比對安全庫存：廠內目前庫存僅剩 2,450 KG，至 9/25 庫存將跌破安全水位 (2,500 KG)。",
        "4. 系統輸出 SOP：建議提前 30 天向國內同業或代理商緊急調料 1,500 KG。"
    ], "系統即時響應", C_PURPLE)

    add_card(s14, 8.8, 1.55, 3.6, 5.35, "3. 跨部門協同處置 (Decision)", [
        "採購處置：採購提前 40 天啟動國內現貨調料 1.5 噸，填補 9/25~10/10 的供應真空期。",
        "生管處置：將非急單排程微調至 10/12 原料進廠後生產，優先保障 ICU 重點訂單。",
        "結果：機台零停線，客戶出貨零延誤，成功化解供應鏈斷鏈危機！"
    ], "危機提早化解", C_EMERALD)

    s14.notes_slide.notes_text_frame.text = (
        "【演講備忘錄 - Slide 14】：情境四是我們最怕的『國外海運延誤』。\n"
        "船期延後 25 天，過去都是到了預定到貨日倉庫沒東西，才雞飛狗跳去借料，甚至機台停線。\n"
        "現在只要採購把新的船期輸進去，系統提前 45 天就亮紅燈，告訴我們 9 月底會斷料，"
        "採購有充裕的時間向國內同行借 1.5 噸先頂著，機台完全不需要停線！"
    )

    # ═════════════════════════════════════════════════════════════════════
    # Slide 15 ~ Slide 18: 資料交換、雙模、效益與 Roadmap
    # ═════════════════════════════════════════════════════════════════

    # Slide 15: 數據中心與 Excel 雙向交換
    s15 = prs.slides.add_slide(blank_layout)
    add_header(s15, "數據中心", "無損資料中心與 Excel / JSON 雙向交換 SOP", "責任角色：全體使用者 / IT 系統維護窗口")
    add_card(s15, 0.8, 1.55, 3.6, 5.35, "標準 Excel 雙向匯出入", [
        "進入『數據中心 ➔ 資料交換與模擬 (DataExchangeView)』。",
        "點擊『下載正式生產空白範本 (.xlsx)』：獲取 9 大工作表標準 Excel 結構。",
        "各部門在專屬工作表填報後，直接拖曳檔案至『智慧匯入區』上傳。"
    ], "標準 Excel 流程", C_BLUE_BRAND)
    add_card(s15, 4.8, 1.55, 3.6, 5.35, "Dry-Run 預檢與安全防呆", [
        "上傳後系統自動執行 Dry-Run 預檢，列出各工作表預計匯入筆數與格式診斷。",
        "主鍵 (PK) 衝突自動執行覆蓋更新 (Upsert)，外鍵 (FK) 無效自動發出防呆警告。",
        "自動校驗 material_class 五層分類與數字格式，確保數據 100% 潔淨。"
    ], "預檢防呆機制", C_EMERALD)
    add_card(s15, 8.8, 1.55, 3.6, 5.35, "全庫備份與資料復原", [
        "點擊『匯出完整系統 JSON 備份檔 (.json)』可一鍵下載資料庫全量快照。",
        "支援定時自動背景備份 (Admin 模式)，可設定每週/每日定時封存。",
        "系統具備高容錯復原能力，任何誤操作皆可一秒匯入備份恢復。"
    ], "全庫備份與安全", C_BLUE_ACCENT)
    s15.notes_slide.notes_text_frame.text = "【演講備忘錄 - Slide 15】：資料中心支援 Excel 無損匯入匯出，大家填寫習慣不變，上傳自帶防呆檢查。"

    # Slide 16: 智慧雙模機制
    s16 = prs.slides.add_slide(blank_layout)
    add_header(s16, "系統創新", "智慧雙模換檔機制 (Demo ↔ Production Mode)", "開箱即用 52 筆全階層示範演練庫 · 匯入真實數據自動無縫換檔")
    add_card(s16, 0.8, 1.55, 3.6, 5.35, "模式一：🎮 示範演練模式 (DEMO)", [
        "開箱預載 52 筆代表性物料數據鏈 (RAW:12, MAT:8, PART:18, COMP:8, SET:6)。",
        "涵蓋 15 組主力模具、良率標準、採購在途訂單，全數據鏈 100% 貫通閉環。",
        "頂部狀態列即時顯示『🎮 示範演練模式 (DEMO)』(天藍標籤)。",
        "用途：新進員工培訓、QCC 成果展示、跨部門推演練習。"
    ], "開箱零冷啟動", C_BLUE_ACCENT)
    add_card(s16, 4.8, 1.55, 3.6, 5.35, "智慧換檔：🟢 正式生產模式 (PROD)", [
        "當各責任單位在資料中心上傳真實生產 Excel / JSON 時：",
        "系統自動清空示範假資料，無縫切換為主體真實生產數據庫！",
        "頂部狀態列自動亮起『🟢 正式生產模式 (PROD)』(翡翠綠標籤)。",
        "徹底防止示範演練資料與真實生產資料交叉污染 (Zero Contamination)。"
    ], "無縫切換正式庫", C_EMERALD)
    add_card(s16, 8.8, 1.55, 3.6, 5.35, "沙盒控制：雙向自由切換", [
        "在資料交換中心隨時提供兩大一鍵換檔按鈕：",
        "👉 🎮 一鍵載入 52 筆示範演練庫：隨時重置回演練沙盒環境。",
        "👉 🧹 一鍵清空資料庫 (切換為純淨空白)：重置為純淨正式空庫供資料填報。",
        "各部門可隨心所欲切換演練與正式狀態，無後顧之憂。"
    ], "沙盒自由切換", C_PURPLE)
    s16.notes_slide.notes_text_frame.text = "【演講備忘錄 - Slide 16】：智慧雙模讓大家練習時有完整的 52 筆示範數據可用，正式上傳時無縫換檔。"

    # Slide 17: 量化效益評估
    s17 = prs.slides.add_slide(blank_layout)
    add_header(s17, "效益量化", "PMS 系統導入量化與質化綜合效益評估", "向管理階層與高階主管匯報之核心商業價值 (ROI)")
    add_card(s17, 0.8, 1.55, 3.6, 5.35, "時間與作業效率提升 95%", [
        "週二雙週出貨審查決策時間：由 120 分鐘大幅縮短至 5 分鐘以內完成。",
        "3 階 MRP 全局推導計算時間：由過去手動拉表 3 天縮短為 0.2 秒即時完成。",
        "跨部門溝通成本降低 80%：消除 Excel 表版本衝突與責任推諉。"
    ], "決策效率躍升", C_EMERALD)
    add_card(s17, 4.8, 1.55, 3.6, 5.35, "庫存資金與呆滯風險歸零", [
        "及時攔截客戶預測下修：採購端自動踩煞車，消除高達 8,000 萬元呆滯料積壓。",
        "倉容分批到貨建議：原料倉平均庫存週轉天數優化 30% 以上，消除走道堆積。",
        "WIP 時序差消除：避免因誤判缺料而重複採購或多開機台。"
    ], "資金佔用大幅下降", C_BLUE_ACCENT)
    add_card(s17, 8.8, 1.55, 3.6, 5.35, "商務信賴度與交付率提升", [
        "客戶預測偏差與供需透明化報告：賦能業務商務談判，化被動責難為主動背書。",
        "準時交付率 (OTD)：預期由 92% 大幅提升至 99.2% 以上。",
        "大幅提升國際一線醫療客戶 (ICU / MDX) 對我方製造管理能力之高度信賴。"
    ], "客戶信賴度提升", C_BLUE_BRAND)
    s17.notes_slide.notes_text_frame.text = "【演講備忘錄 - Slide 17】：這是向總經理與各處主管報告的核心效益：開會變快、呆料歸零、客戶信賴提升！"

    # Slide 18: Roadmap 與 Q&A
    s18 = prs.slides.add_slide(blank_layout)
    add_header(s18, "行動推進", "系統導入 Roadmap、跨部門分工與 Q&A", "按部就班穩健推進 · 預計明年初正式全面上線")
    add_card(s18, 0.8, 1.55, 3.6, 5.35, "Phase 1: 主檔填報 (8~9月)", [
        "業務部：整理 ICU / MDX 年度 Forecast 與近期正式訂單。",
        "工程/品保：審查製品 BOM、模具克重、損耗率與全檢良率。",
        "製造部：確認 15 組主力模具妥善穴數與週期標準值。",
        "資材部：確認供應商交期、MOQ 與實體倉容限制。"
    ], "現階段任務 (進行中)", C_BLUE_ACCENT)
    add_card(s18, 4.8, 1.55, 3.6, 5.35, "Phase 2: 雙軌試跑 (10~11月)", [
        "雙軌並行 (Parallel Run)：每週二出貨會議同步使用舊 Excel 表與 PMS 系統驗證比對。",
        "收集現場使用者回饋：進行介面微調與操作體驗優化。",
        "完成與頂新 ERP 資料庫介面之拋轉對接與防呆測試。"
    ], "第二階段 (即將啟動)", C_AMBER)
    add_card(s18, 8.8, 1.55, 3.6, 5.35, "Phase 3: 正式上線 (明年初)", [
        "1 月 1 號/4 號正式全面切換至【正式生產模式】。",
        "全面停用舊版手工表格，以 PMS 作為全廠單一事實來源 (SSOT)。",
        "支援窗口：Wesley Chang (系統主導) / 各部門種子同仁。",
        "現場 Q&A 交流時間，歡迎各位長官與夥伴提問！"
    ], "正式全面上線", C_EMERALD)
    s18.notes_slide.notes_text_frame.text = (
        "【演講備忘錄 - Slide 18】：最後是推進時程，預計明年初全面上線。\n"
        "感謝各位長官與跨部門同仁的全力支持！接下來是現場 Q&A 時間，請大家踴躍提問與交流，謝謝大家！"
    )

    output_path = "docs/PMS_跨部門系統操作與決策協同培訓簡報.pptx"
    prs.save(output_path)
    print(f"[SUCCESS] 18 頁超詳細簡報（含 4 大情境實戰模擬）已成功生成至: {os.path.abspath(output_path)}")
    return output_path

if __name__ == "__main__":
    create_presentation()
